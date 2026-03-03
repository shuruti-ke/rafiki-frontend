import logging
import re
from pathlib import Path
from sqlalchemy.orm import Session
from app.models.document import Document, DocumentChunk

logger = logging.getLogger(__name__)


def extract_text_from_bytes(content_bytes: bytes, mime_type: str) -> str:
    """Extract text from in-memory file bytes (for R2-stored documents)."""
    if mime_type == "application/pdf":
        return _extract_pdf_bytes(content_bytes)
    elif mime_type in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"):
        return _extract_docx_bytes(content_bytes)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx_bytes(content_bytes)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx_bytes(content_bytes)
    elif mime_type.startswith("text/"):
        return content_bytes.decode("utf-8", errors="replace")
    else:
        logger.warning("Unsupported mime type for text extraction: %s", mime_type)
        return ""


def _extract_pdf_bytes(data: bytes) -> str:
    """
    Extract text from PDF bytes.
    Primary:  pdfminer.six  — handles modern/complex PDFs reliably.
    Fallback: PyPDF2        — kept for compatibility.
    """
    # --- Primary: pdfminer.six ---
    try:
        import io
        from pdfminer.high_level import extract_text as pdfminer_extract_text
        from pdfminer.layout import LAParams

        text = pdfminer_extract_text(
            io.BytesIO(data),
            laparams=LAParams(line_margin=0.5, word_margin=0.1),
        )
        if text and text.strip():
            logger.info("PDF extracted via pdfminer.six (%d chars)", len(text))
            return text.strip()
        logger.warning("pdfminer.six returned empty text — trying PyPDF2 fallback")
    except ImportError:
        logger.warning("pdfminer.six not installed. Run: pip install pdfminer.six")
    except Exception as e:
        logger.warning("pdfminer.six failed (%s) — trying PyPDF2 fallback", e)

    # --- Fallback: PyPDF2 ---
    try:
        import io
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            logger.warning("PDF is encrypted — cannot extract text without password")
            return ""

        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        result = "\n\n".join(pages)
        if result.strip():
            logger.info("PDF extracted via PyPDF2 fallback (%d chars)", len(result))
            return result
        logger.warning("PyPDF2 also returned empty — PDF may be scanned/image-based")
        return ""
    except Exception as e:
        logger.error("PDF bytes extraction failed entirely: %s", e)
        return ""


def _extract_docx_bytes(data: bytes) -> str:
    try:
        import io
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.error("DOCX bytes extraction failed: %s", e)
        return ""


def _extract_xlsx_bytes(data: bytes) -> str:
    try:
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
        wb.close()
        return "\n".join(rows)
    except Exception as e:
        logger.error("XLSX bytes extraction failed: %s", e)
        return ""


def _extract_pptx_bytes(data: bytes) -> str:
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error("PPTX bytes extraction failed: %s", e)
        return ""


def extract_text(file_path: str, mime_type: str) -> str:
    """Extract text from files on the local filesystem (fallback path)."""
    full_path = Path(__file__).parent.parent.parent / file_path

    if mime_type == "application/pdf":
        return _extract_pdf_bytes(full_path.read_bytes())
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx_bytes(full_path.read_bytes())
    elif mime_type.startswith("text/"):
        return full_path.read_text(encoding="utf-8", errors="replace")
    else:
        logger.warning("Unsupported mime type for text extraction: %s", mime_type)
        return ""


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    """Split text into chunks on paragraph/sentence boundaries."""
    if not text.strip():
        return []

    paragraphs = re.split(r"\n\s*\n", text)
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk = f"{current_chunk}\n\n{para}" if current_chunk else para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
                words = current_chunk.split()
                overlap_words = words[-overlap // 4:] if len(words) > overlap // 4 else []
                current_chunk = " ".join(overlap_words) + "\n\n" + para if overlap_words else para
            else:
                sentences = re.split(r"(?<=[.!?])\s+", para)
                for sent in sentences:
                    if len(current_chunk) + len(sent) + 1 <= chunk_size:
                        current_chunk = f"{current_chunk} {sent}" if current_chunk else sent
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sent

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def _estimate_tokens(text: str) -> int:
    return len(text) // 4


def index_document(db: Session, document: Document):
    """Extract text, chunk it, and store chunks for full-text search."""
    text = ""
    try:
        from app.services.file_storage import download_file_bytes
        file_bytes = download_file_bytes(document.file_path)
        text = extract_text_from_bytes(file_bytes, document.mime_type)
    except Exception as e:
        logger.debug("R2 download failed, trying local: %s", e)
        text = extract_text(document.file_path, document.mime_type)

    if not text:
        logger.warning("No text extracted from document %s", document.id)
        return

    db.query(DocumentChunk).filter(DocumentChunk.document_id == document.id).delete()

    chunks = chunk_text(text)
    for i, chunk_content in enumerate(chunks):
        chunk = DocumentChunk(
            document_id=document.id,
            org_id=document.org_id,
            chunk_index=i,
            content=chunk_content,
            token_count=_estimate_tokens(chunk_content),
        )
        db.add(chunk)

    document.is_indexed = True
    db.commit()
    logger.info("Indexed document %d: %d chunks", document.id, len(chunks))

